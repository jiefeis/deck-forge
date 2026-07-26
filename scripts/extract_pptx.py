#!/usr/bin/env python3
"""
Extract a PowerPoint file (.pptx) as source material for deck generation.

This extractor is intentionally lossy. Its JSON output is useful for reading
and rebuilding a deck, but it is never a representation for native PPTX edits
or round-tripping. Native edits must operate on the original PPTX package.

The output remains a list of slide dictionaries for backward compatibility.
Each slide records its source order, stable slide id, visibility, layout, canvas
size, text, tables, images, and notes. Shape geometry uses EMUs.

Usage:
    python extract_pptx.py <input.pptx> [output_dir] [--visible-only]

Requires: pip install python-pptx
"""

import argparse
import json
import math
import os
import sys

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    sys.exit("extract_pptx: python-pptx required. Run: pip install python-pptx")


SOURCE_MATERIAL_NOTICE = (
    "Source-material extraction only; do not use this JSON for native PPTX "
    "editing or round-tripping."
)

# Affine transform (a, b, c, d, e, f):
# x' = a*x + c*y + e; y' = b*x + d*y + f.
IDENTITY_TRANSFORM = (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)


def _compose(outer, inner):
    """Return the affine transform outer(inner(point))."""
    oa, ob, oc, od, oe, of = outer
    ia, ib, ic, id_, ie, if_ = inner
    return (
        oa * ia + oc * ib,
        ob * ia + od * ib,
        oa * ic + oc * id_,
        ob * ic + od * id_,
        oa * ie + oc * if_ + oe,
        ob * ie + od * if_ + of,
    )


def _apply_transform(transform, x, y):
    a, b, c, d, e, f = transform
    return a * x + c * y + e, b * x + d * y + f


def _rotation_about(angle_degrees, center_x, center_y):
    """Return a clockwise rotation in PowerPoint's y-down coordinate space."""
    angle = math.radians(angle_degrees)
    cosine = math.cos(angle)
    sine = math.sin(angle)
    return (
        cosine,
        sine,
        -sine,
        cosine,
        center_x - cosine * center_x + sine * center_y,
        center_y - sine * center_x - cosine * center_y,
    )


def _flip_about(flip_h, flip_v, center_x, center_y):
    scale_x = -1.0 if flip_h else 1.0
    scale_y = -1.0 if flip_v else 1.0
    return (
        scale_x,
        0.0,
        0.0,
        scale_y,
        center_x * (1.0 - scale_x),
        center_y * (1.0 - scale_y),
    )


def _group_child_transform(group_shape):
    """Map a group's child coordinate system into its parent's coordinates."""
    xfrm = group_shape._element.grpSpPr.xfrm
    if (
        xfrm is None
        or xfrm.off is None
        or xfrm.ext is None
        or xfrm.chOff is None
        or xfrm.chExt is None
        or not xfrm.chExt.cx
        or not xfrm.chExt.cy
    ):
        return IDENTITY_TRANSFORM

    off_x = int(xfrm.off.x)
    off_y = int(xfrm.off.y)
    ext_x = int(xfrm.ext.cx)
    ext_y = int(xfrm.ext.cy)
    child_off_x = int(xfrm.chOff.x)
    child_off_y = int(xfrm.chOff.y)
    scale_x = ext_x / int(xfrm.chExt.cx)
    scale_y = ext_y / int(xfrm.chExt.cy)

    scale_and_offset = (
        scale_x,
        0.0,
        0.0,
        scale_y,
        off_x - scale_x * child_off_x,
        off_y - scale_y * child_off_y,
    )
    center_x = off_x + ext_x / 2.0
    center_y = off_y + ext_y / 2.0
    flipped = _compose(
        _flip_about(xfrm.flipH, xfrm.flipV, center_x, center_y),
        scale_and_offset,
    )
    return _compose(
        _rotation_about(float(xfrm.rot or 0.0), center_x, center_y),
        flipped,
    )


def _shape_bbox(shape, transform):
    """Return the shape bounds in slide coordinates as integer EMUs."""
    left = int(shape.left or 0)
    top = int(shape.top or 0)
    width = int(shape.width or 0)
    height = int(shape.height or 0)
    corners = (
        _apply_transform(transform, left, top),
        _apply_transform(transform, left + width, top),
        _apply_transform(transform, left, top + height),
        _apply_transform(transform, left + width, top + height),
    )
    xs = [point[0] for point in corners]
    ys = [point[1] for point in corners]
    bbox_left = round(min(xs))
    bbox_top = round(min(ys))
    return {
        "left": bbox_left,
        "top": bbox_top,
        "width": round(max(xs)) - bbox_left,
        "height": round(max(ys)) - bbox_top,
    }


def walk_shape_records(
    shapes, parent_transform=IDENTITY_TRANSFORM, parent_z_path=()
):
    """Yield leaf shapes with slide-space bounds and exact z-order paths.

    ``z_order`` is the top-level drawing order. ``z_order_path`` retains the
    full local order for shapes nested in groups, for example ``[3, 1]``.
    """
    for local_z_order, shape in enumerate(shapes):
        z_path = parent_z_path + (local_z_order,)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            child_transform = _compose(
                parent_transform, _group_child_transform(shape)
            )
            yield from walk_shape_records(
                shape.shapes, child_transform, z_path
            )
            continue
        yield {
            "shape": shape,
            "bbox": _shape_bbox(shape, parent_transform),
            "z_order": z_path[0],
            "z_order_path": list(z_path),
        }


def walk_shapes(shapes):
    """Yield leaf shapes depth-first, expanding group shapes.

    Kept as a compatibility helper for callers that only need shape objects.
    """
    for record in walk_shape_records(shapes):
        yield record["shape"]


def _shape_metadata(record):
    shape = record["shape"]
    shape_type = shape.shape_type
    shape_type_name = getattr(shape_type, "name", None)
    if shape_type_name is None:
        shape_type_name = str(shape_type).split(" (", 1)[0]
    return {
        "shape_name": shape.name,
        "shape_id": int(shape.shape_id),
        "shape_type": shape_type_name,
        "bbox": record["bbox"],
        "z_order": record["z_order"],
        "z_order_path": record["z_order_path"],
    }


def is_slide_hidden(slide):
    """Return whether the source slide is explicitly marked not to show."""
    show_value = slide._element.get("show")
    if show_value is None:
        return False
    return str(show_value).strip().lower() in {"0", "false", "off", "no"}


def get_image(shape, slide_num):
    """Return the shape's embedded image object, or None.

    Covers plain pictures and pictures placed inside placeholders. Linked
    (non-embedded) pictures raise on access; warn and continue instead of
    crashing the whole extraction.
    """
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE and not shape.is_placeholder:
        return None
    try:
        return shape.image
    except AttributeError:
        return None  # placeholder without an inserted picture
    except Exception as exc:
        print(
            f"  warning: slide {slide_num}: skipping unreadable image "
            f"(linked, not embedded?): {exc}"
        )
        return None


def _extract_pptx(file_path, output_dir=".", visible_only=False):
    prs = Presentation(file_path)
    slides_data = []
    hidden_slide_numbers = []
    canvas = {
        "width": int(prs.slide_width),
        "height": int(prs.slide_height),
        "unit": "EMU",
    }

    # Create assets directory for extracted images.
    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    # python-pptx iterates the p:sldIdLst relationship order, which is the
    # presentation's true slide order rather than slideN.xml filename order.
    for slide_num, slide in enumerate(prs.slides, start=1):
        hidden = is_slide_hidden(slide)
        if hidden:
            hidden_slide_numbers.append(slide_num)
            if visible_only:
                continue

        layout_name = slide.slide_layout.name or ""
        slide_data = {
            "number": slide_num,
            "slide_id": int(slide.slide_id),
            "hidden": hidden,
            "layout": layout_name,
            "canvas": dict(canvas),
            "title": "",
            "title_shape": None,
            "content": [],
            "images": [],
            "notes": "",
        }
        title_shape = slide.shapes.title
        title_shape_id = (
            int(title_shape.shape_id) if title_shape is not None else None
        )

        for record in walk_shape_records(slide.shapes):
            shape = record["shape"]
            metadata = _shape_metadata(record)

            # Extract text while retaining the historical title/body split.
            if shape.has_text_frame:
                if title_shape_id is not None and shape.shape_id == title_shape_id:
                    slide_data["title"] = shape.text
                    slide_data["title_shape"] = metadata
                elif shape.text.strip():
                    slide_data["content"].append(
                        {
                            "type": "text",
                            "content": shape.text,
                            **metadata,
                        }
                    )

            # Extract table cell text.
            if shape.has_table:
                rows = [
                    [cell.text for cell in row.cells]
                    for row in shape.table.rows
                ]
                slide_data["content"].append(
                    {"type": "table", "rows": rows, **metadata}
                )

            # Extract images (embedded pictures, including picture placeholders).
            image = get_image(shape, slide_num)
            if image is None:
                continue
            try:
                image_bytes = image.blob
                image_ext = image.ext
            except Exception as exc:
                print(
                    f"  warning: slide {slide_num}: could not read "
                    f"image data: {exc}"
                )
                continue
            image_name = (
                f"slide{slide_num}_img{len(slide_data['images']) + 1}."
                f"{image_ext}"
            )
            image_path = os.path.join(assets_dir, image_name)

            with open(image_path, "wb") as image_file:
                image_file.write(image_bytes)

            bbox = metadata["bbox"]
            slide_data["images"].append(
                {
                    "path": f"assets/{image_name}",
                    # Keep the historical width/height fields and add position.
                    "left": bbox["left"],
                    "top": bbox["top"],
                    "width": bbox["width"],
                    "height": bbox["height"],
                    **metadata,
                }
            )

        # Extract speaker notes.
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            slide_data["notes"] = notes_frame.text

        slides_data.append(slide_data)

    summary = {
        "source_slide_count": len(prs.slides),
        "extracted_slide_count": len(slides_data),
        "hidden_slide_numbers": hidden_slide_numbers,
    }
    return slides_data, summary


def extract_pptx(file_path, output_dir=".", visible_only=False):
    """Extract lossy source material; never use the result for native edits.

    Returns the historical list of slide dictionaries. By default hidden slides
    are retained and identified with ``hidden: true``. Set ``visible_only`` to
    omit them while preserving each remaining slide's original ``number``.
    """
    slides_data, _summary = _extract_pptx(
        file_path, output_dir, visible_only=visible_only
    )
    return slides_data


def _slide_number_summary(slide_numbers):
    if not slide_numbers:
        return ""
    label = "source slide" if len(slide_numbers) == 1 else "source slides"
    return f" ({label} {', '.join(str(n) for n in slide_numbers)})"


def main():
    # Keep console output safe on non-UTF-8 Windows pipes (GBK/cp1252).
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(errors="replace")

    ap = argparse.ArgumentParser(
        description=(
            "Extract a .pptx as lossy source material for deck generation. "
            "This output is not suitable for native PPTX editing."
        ),
        epilog=SOURCE_MATERIAL_NOTICE,
    )
    ap.add_argument("input_pptx", help="path to the .pptx file")
    ap.add_argument(
        "output_dir",
        nargs="?",
        default=".",
        help=(
            "directory for extracted-slides.json and assets/ "
            "(default: current dir)"
        ),
    )
    ap.add_argument(
        "--visible-only",
        action="store_true",
        help=(
            "omit hidden slides; retained slides keep their original source "
            "slide numbers"
        ),
    )
    args = ap.parse_args()

    slides, summary = _extract_pptx(
        args.input_pptx,
        args.output_dir,
        visible_only=args.visible_only,
    )

    # Write extracted data as JSON.
    output_path = os.path.join(args.output_dir, "extracted-slides.json")
    with open(output_path, "w", encoding="utf-8") as output_file:
        json.dump(slides, output_file, indent=2, ensure_ascii=False)

    print(SOURCE_MATERIAL_NOTICE)
    print(f"Extracted {len(slides)} slides to {output_path}")
    hidden_numbers = summary["hidden_slide_numbers"]
    hidden_detail = _slide_number_summary(hidden_numbers)
    if args.visible_only:
        print(
            "Hidden slides omitted by --visible-only: "
            f"{len(hidden_numbers)}{hidden_detail}."
        )
    elif hidden_numbers:
        print(
            f"Hidden slides retained: {len(hidden_numbers)}{hidden_detail}. "
            "Use --visible-only to omit them."
        )
    else:
        print("Hidden slides retained: 0.")

    for slide in slides:
        img_count = len(slide["images"])
        visibility = " [hidden]" if slide["hidden"] else ""
        print(
            f"  Slide {slide['number']}{visibility}: "
            f"{slide['title'] or '(no title)'} - {img_count} image(s)"
        )


if __name__ == "__main__":
    main()
