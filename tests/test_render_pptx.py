#!/usr/bin/env python3
"""Windows integration test for scripts/render_pptx.ps1."""

from __future__ import annotations

import hashlib
import json
import os
import shutil
import subprocess
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parent.parent
SCRIPT = ROOT / "scripts" / "render_pptx.ps1"


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


class RenderPptxTests(unittest.TestCase):
    @unittest.skipUnless(os.name == "nt" and shutil.which("powershell.exe"),
                         "PowerPoint/WPS COM rendering is Windows-only")
    def test_visible_and_hidden_render_without_source_mutation(self) -> None:
        temp = Path(tempfile.mkdtemp(prefix="deckforge-render-test-"))
        try:
            pptx = temp / "source.pptx"
            prs = Presentation()
            visible = prs.slides.add_slide(prs.slide_layouts[1])
            visible.shapes.title.text = "VISIBLE"
            hidden = prs.slides.add_slide(prs.slide_layouts[1])
            hidden.shapes.title.text = "HIDDEN"
            hidden._element.set("show", "0")
            prs.save(pptx)
            before = sha256(pptx)

            visible_dir = temp / "visible"
            all_dir = temp / "all"
            base = [
                "powershell.exe", "-NoProfile", "-ExecutionPolicy", "Bypass",
                "-File", str(SCRIPT), "-InputPptx", str(pptx),
            ]
            first = subprocess.run(
                base + ["-OutputDir", str(visible_dir), "-Engine", "auto"],
                capture_output=True, text=True, timeout=90,
            )
            if first.returncode and "No supported presentation engine" in (
                first.stdout + first.stderr
            ):
                self.skipTest("PowerPoint/WPS COM engine is not registered")
            self.assertEqual(first.returncode, 0, first.stdout + first.stderr)
            second = subprocess.run(
                base + ["-OutputDir", str(all_dir), "-Engine", "auto",
                        "-IncludeHidden"],
                capture_output=True, text=True, timeout=90,
            )
            self.assertEqual(second.returncode, 0, second.stdout + second.stderr)

            self.assertEqual(before, sha256(pptx))
            self.assertEqual(len(list(visible_dir.glob("slide-*.png"))), 1)
            self.assertEqual(len(list(all_dir.glob("slide-*.png"))), 2)
            manifest = json.loads(
                (visible_dir / "render-manifest.json").read_text(encoding="utf-8-sig")
            )
            self.assertEqual(manifest["slide_count"], 2)
            self.assertEqual(manifest["hidden_physical_indices"], [2])
            self.assertFalse(manifest["include_hidden"])
        finally:
            shutil.rmtree(temp, ignore_errors=True)


if __name__ == "__main__":
    unittest.main(verbosity=2)
