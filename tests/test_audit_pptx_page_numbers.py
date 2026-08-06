#!/usr/bin/env python3
"""Regression tests for the strict PPTX page-number audit."""

from __future__ import annotations

import os
import subprocess
import sys
import tempfile
import unittest
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


sys.dont_write_bytecode = True

ROOT = Path(__file__).resolve().parent.parent
SCRIPT = ROOT / "scripts" / "audit_pptx_page_numbers.py"

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
for prefix, uri in NS.items():
    ET.register_namespace(prefix, uri)


def _is_native_number_shape(shape: ET.Element) -> bool:
    placeholders = shape.findall(".//p:ph", NS)
    fields = shape.findall(".//a:fld", NS)
    return any(
        (node.get("type") or "").lower() == "sldnum"
        for node in placeholders
    ) or any(
        (node.get("type") or "").lower() == "slidenum" for node in fields
    )


def _strip_native_number_shapes(path: Path, *, keep_master: bool = False) -> None:
    """Remove template-native number placeholders while preserving the package."""
    rewritten = path.with_name(path.stem + ".rewritten.pptx")
    with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(
        rewritten, "w"
    ) as target:
        for info in source.infolist():
            data = source.read(info.filename)
            is_number_part = (
                info.filename.startswith("ppt/slides/")
                or info.filename.startswith("ppt/slideLayouts/")
                or info.filename.startswith("ppt/slideMasters/")
            ) and info.filename.endswith(".xml")
            preserve_part = keep_master and info.filename.startswith(
                "ppt/slideMasters/"
            )
            if is_number_part and not preserve_part:
                root = ET.fromstring(data)
                for parent in root.iter():
                    for child in list(parent):
                        if child.tag == f"{{{NS['p']}}}sp" and _is_native_number_shape(
                            child
                        ):
                            parent.remove(child)
                data = ET.tostring(
                    root, encoding="utf-8", xml_declaration=True
                )
            target.writestr(info, data)
    rewritten.replace(path)


def _add_direct_marker(
    presentation: Presentation,
    slide,
    value: int,
    *,
    name: str,
    size_pt: int,
    offset: int = 0,
) -> None:
    marker = slide.shapes.add_textbox(
        presentation.slide_width - Inches(0.8 + offset * 0.8),
        presentation.slide_height - Inches(0.45),
        Inches(0.6),
        Inches(0.3),
    )
    marker.name = name
    marker.text_frame.clear()
    run = marker.text_frame.paragraphs[0].add_run()
    run.text = str(value)
    run.font.size = Pt(size_pt)


def _make_deck(
    path: Path,
    values: list[int | None],
    *,
    hidden: set[int] | None = None,
    duplicates: dict[int, list[int]] | None = None,
    name: str = "page_number",
    size_pt: int = 11,
    keep_native_master: bool = False,
) -> Path:
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    hidden = hidden or set()
    duplicates = duplicates or {}

    for physical_number, value in enumerate(values, 1):
        slide = presentation.slides.add_slide(blank)
        if physical_number in hidden:
            slide._element.set("show", "0")
        if value is not None:
            _add_direct_marker(
                presentation,
                slide,
                value,
                name=name,
                size_pt=size_pt,
            )
        for offset, duplicate_value in enumerate(
            duplicates.get(physical_number, []), 1
        ):
            _add_direct_marker(
                presentation,
                slide,
                duplicate_value,
                name=name,
                size_pt=size_pt,
                offset=offset,
            )

    presentation.save(path)
    _strip_native_number_shapes(path, keep_master=keep_native_master)
    return path


class AuditPptxPageNumbersTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp_dir = tempfile.TemporaryDirectory(
            prefix="deckforge-page-number-test-"
        )
        self.temp = Path(self.temp_dir.name)

    def tearDown(self) -> None:
        self.temp_dir.cleanup()

    def run_audit(self, path: Path, *args: str) -> subprocess.CompletedProcess[str]:
        env = os.environ.copy()
        env["PYTHONUTF8"] = "1"
        return subprocess.run(
            [sys.executable, "-X", "utf8", str(SCRIPT), *args, str(path)],
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            env=env,
            check=False,
        )

    def assert_ok(self, result: subprocess.CompletedProcess[str]) -> None:
        self.assertEqual(
            result.returncode,
            0,
            msg=f"stdout:\n{result.stdout}\nstderr:\n{result.stderr}",
        )
        self.assertIn("status: OK", result.stdout)

    def assert_failed(self, result: subprocess.CompletedProcess[str]) -> None:
        self.assertNotEqual(
            result.returncode,
            0,
            msg=f"stdout:\n{result.stdout}\nstderr:\n{result.stderr}",
        )
        self.assertIn("status: FAIL", result.stdout)

    def test_auto_accepts_valid_direct_numbers(self) -> None:
        path = _make_deck(
            self.temp / "\u9875\u7801\u5ba1\u8ba1.pptx", [1, 2, 3]
        )

        result = self.run_audit(path)

        self.assert_ok(result)
        self.assertIn("mode: direct (requested auto)", result.stdout)
        self.assertIn("direct page-number shapes: 3", result.stdout)

    def test_auto_rejects_duplicate_and_missing_direct_numbers(self) -> None:
        duplicate = _make_deck(
            self.temp / "duplicate.pptx",
            [1, 2, 3],
            duplicates={2: [2]},
        )
        missing = _make_deck(self.temp / "missing.pptx", [1, None, 3])

        duplicate_result = self.run_audit(duplicate)
        missing_result = self.run_audit(missing)

        self.assert_failed(duplicate_result)
        self.assertIn("slide 2: duplicate direct page numbers", duplicate_result.stdout)
        self.assert_failed(missing_result)
        self.assertIn("slide 2: missing direct page number", missing_result.stdout)

    def test_auto_rejects_direct_and_inherited_native_conflict(self) -> None:
        path = _make_deck(
            self.temp / "inherited-conflict.pptx",
            [1, 2],
            keep_native_master=True,
        )

        result = self.run_audit(path)

        self.assert_failed(result)
        self.assertIn("inherited/native page-number sources coexist", result.stdout)

    def test_hidden_slides_follow_physical_or_visible_numbering(self) -> None:
        physical = _make_deck(
            self.temp / "physical.pptx", [1, 2, 3], hidden={2}
        )
        visible = _make_deck(
            self.temp / "visible.pptx", [1, None, 2], hidden={2}
        )

        physical_result = self.run_audit(
            physical, "--numbering", "physical"
        )
        visible_result = self.run_audit(visible, "--numbering", "visible")
        wrong_physical = self.run_audit(visible, "--numbering", "physical")
        wrong_visible = self.run_audit(physical, "--numbering", "visible")

        self.assert_ok(physical_result)
        self.assertIn("hidden 1, audited 3", physical_result.stdout)
        self.assert_ok(visible_result)
        self.assertIn("hidden 1, audited 2", visible_result.stdout)
        self.assert_failed(wrong_physical)
        self.assertIn("slide 2: missing direct page number", wrong_physical.stdout)
        self.assert_failed(wrong_visible)
        self.assertIn("slide 3: expected page number 2", wrong_visible.stdout)

    def test_numbering_none_skips_presence_and_sequence_only(self) -> None:
        path = _make_deck(self.temp / "unsequenced.pptx", [10, None, 99])

        strict_result = self.run_audit(path)
        relaxed_result = self.run_audit(path, "--numbering", "none")

        self.assert_failed(strict_result)
        self.assert_ok(relaxed_result)

    def test_mode_native_accepts_one_reachable_native_source(self) -> None:
        path = _make_deck(
            self.temp / "native.pptx",
            [None, None],
            keep_native_master=True,
        )

        native_result = self.run_audit(path, "--mode", "native")
        auto_result = self.run_audit(path)
        direct_result = self.run_audit(path, "--mode", "direct")

        self.assert_ok(native_result)
        self.assert_ok(auto_result)
        self.assertIn("mode: native (requested auto)", auto_result.stdout)
        self.assert_failed(direct_result)
        self.assertIn("direct mode found inherited/native", direct_result.stdout)

    def test_expected_name_and_size_are_optional_strict_checks(self) -> None:
        path = _make_deck(
            self.temp / "metadata.pptx",
            [1, 2],
            name="page_number",
            size_pt=11,
        )

        valid = self.run_audit(
            path,
            "--mode",
            "direct",
            "--expect-name",
            "page_number",
            "--expect-size",
            "1100",
        )
        wrong_name = self.run_audit(
            path, "--expect-name", "different_name"
        )
        wrong_size = self.run_audit(path, "--expect-size", "1200")

        self.assert_ok(valid)
        self.assert_failed(wrong_name)
        self.assertIn("expected page-number name", wrong_name.stdout)
        self.assert_failed(wrong_size)
        self.assertIn("expected page-number size 1200", wrong_size.stdout)


    def test_relationship_target_fragment_is_stripped(self) -> None:
        """A '#fragment' on an internal Target must resolve, not crash the audit.

        resolve_target strips the fragment so this script agrees with the other
        audit scripts; without it the lookup raised a raw zip KeyError and the
        whole deck reported 'cannot audit PPTX'.
        """
        path = _make_deck(self.temp / "fragment.pptx", [1, 2, 3])
        rels = "ppt/_rels/presentation.xml.rels"
        with zipfile.ZipFile(path) as archive:
            members = {n: archive.read(n) for n in archive.namelist()}
        patched = members[rels].decode("utf-8").replace(
            'Target="slides/slide1.xml"', 'Target="slides/slide1.xml#anchor"', 1
        )
        self.assertIn("#anchor", patched, "fixture did not contain the expected Target")
        members[rels] = patched.encode("utf-8")
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
            for name, data in members.items():
                archive.writestr(name, data)

        result = self.run_audit(path)
        self.assertNotIn("cannot audit PPTX", result.stdout)
        self.assert_ok(result)


if __name__ == "__main__":
    unittest.main(verbosity=2)
