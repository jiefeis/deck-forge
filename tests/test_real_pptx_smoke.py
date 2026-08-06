#!/usr/bin/env python3
"""Smoke test: run the audit/transplant public APIs on a real python-pptx deck."""

from __future__ import annotations

import importlib.util
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parent.parent
SCRIPTS = ROOT / "scripts"
SPEC = importlib.util.spec_from_file_location(
    "audit_pptx_structure", SCRIPTS / "audit_pptx_structure.py"
)
STRUCTURE = importlib.util.module_from_spec(SPEC)
SPEC.loader.exec_module(STRUCTURE)


def build_deck(path: Path) -> None:
    prs = Presentation()
    for i in range(4):
        # Blank layout with plain text boxes: transplant's safety gate rejects
        # placeholder shapes (ph depends on layout/theme semantics), and this
        # smoke test wants the no-op transplant path covered too.
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = f"Slide {i + 1}"
        if i == 2:
            slide._element.set("show", "0")  # hidden, as in test_render_pptx
    prs.save(path)


class RealPptxSmokeTests(unittest.TestCase):
    def test_real_package_round_trip(self) -> None:
        with tempfile.TemporaryDirectory(prefix="deckforge-smoke-") as temp:
            deck = Path(temp) / "deck.pptx"
            build_deck(deck)

            manifest = STRUCTURE.build_manifest(deck)
            self.assertEqual(manifest["integrity_warnings"], [])
            self.assertEqual(
                [slide["hidden"] for slide in manifest["slides"]],
                [False, False, True, False],
            )

            report = STRUCTURE.compare_pptx(deck, deck)
            self.assertEqual(report["differences"], [])
            self.assertTrue(report["ok"])

            def run_cli(script: str, *args: str) -> None:
                result = subprocess.run(
                    [sys.executable, "-B", str(SCRIPTS / script), *args],
                    capture_output=True, text=True, encoding="utf-8",
                    errors="replace", timeout=120,
                )
                self.assertEqual(
                    result.returncode, 0,
                    f"{script}: {result.stdout}{result.stderr}",
                )

            run_cli("audit_pptx_structure.py", "manifest", str(deck))
            run_cli("audit_pptx_page_numbers.py", str(deck))
            run_cli("audit_pptx_typography.py", str(deck))

            # No-op transplant: candidate is a byte-identical baseline copy, so
            # titles corroborate slide identity and the guards accept page 2.
            candidate = Path(temp) / "candidate.pptx"
            rebased = Path(temp) / "rebased.pptx"
            candidate.write_bytes(deck.read_bytes())
            run_cli(
                "transplant_pptx_slides.py", "--pages", "2",
                str(deck), str(candidate), str(rebased),
            )
            self.assertTrue(rebased.is_file())


if __name__ == "__main__":
    unittest.main(verbosity=2)
