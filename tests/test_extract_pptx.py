#!/usr/bin/env python3
"""Regression tests for scripts/extract_pptx.py.

Run standalone:                    python tests/test_extract_pptx.py
Or under pytest if installed:      python -m pytest tests/test_extract_pptx.py
"""

import importlib.util
import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


sys.dont_write_bytecode = True

ROOT = Path(__file__).resolve().parent.parent
SCRIPT = ROOT / "scripts" / "extract_pptx.py"
SPEC = importlib.util.spec_from_file_location("extract_pptx", SCRIPT)
EXTRACT = importlib.util.module_from_spec(SPEC)
SPEC.loader.exec_module(EXTRACT)


class ExtractPptxTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory(prefix="deckforge-pptx-test-")
        self.root = Path(self.temp_dir.name)
        self.pptx_path = self.root / "source.pptx"
        self.expected = self._build_fixture(self.pptx_path)

    def tearDown(self):
        self.temp_dir.cleanup()

    def _build_fixture(self, path):
        prs = Presentation()
        prs.slide_width = Inches(12)
        prs.slide_height = Inches(7)

        first = prs.slides.add_slide(prs.slide_layouts[1])
        first.shapes.title.text = "First title"
        first.shapes.title.name = "First Title Shape"
        body = first.placeholders[1]
        body.text = "First body"
        body.name = "First Body Shape"

        hidden = prs.slides.add_slide(prs.slide_layouts[5])
        hidden.shapes.title.text = "Hidden title"
        hidden.shapes.title.name = "Hidden Title Shape"
        hidden._element.set("show", "0")

        geometry = prs.slides.add_slide(prs.slide_layouts[6])
        text_shape = geometry.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(1)
        )
        text_shape.name = "Positioned Text"
        text_shape.text = "Positioned body"

        table_shape = geometry.shapes.add_table(
            1, 2, Inches(3), Inches(1), Inches(2), Inches(1)
        )
        table_shape.name = "Positioned Table"
        table_shape.table.cell(0, 0).text = "A"
        table_shape.table.cell(0, 1).text = "B"

        image_path = self.root / "pixel.png"
        Image.new("RGB", (4, 3), (220, 30, 30)).save(image_path)
        picture = geometry.shapes.add_picture(
            str(image_path),
            Inches(6),
            Inches(1),
            width=Inches(2),
            height=Inches(1),
        )
        picture.name = "Positioned Picture"

        group = geometry.shapes.add_group_shape()
        group.name = "Scaled Group"
        grouped_text = group.shapes.add_textbox(
            Inches(1), Inches(4), Inches(2), Inches(1)
        )
        grouped_text.name = "Grouped Text"
        grouped_text.text = "Grouped body"
        grouped_picture = group.shapes.add_picture(
            str(image_path),
            Inches(4),
            Inches(4),
            width=Inches(1),
            height=Inches(1),
        )
        grouped_picture.name = "Grouped Picture"
        # Scale and move the group after its child coordinate space is fixed.
        group.left = Inches(7)
        group.top = Inches(4)
        group.width = Inches(4)
        group.height = Inches(2)

        ids = {
            "first": first.slide_id,
            "hidden": hidden.slide_id,
            "geometry": geometry.slide_id,
        }
        layouts = {
            first.slide_id: first.slide_layout.name,
            hidden.slide_id: hidden.slide_layout.name,
            geometry.slide_id: geometry.slide_layout.name,
        }

        # Deliberately make the true presentation order differ from creation
        # order (and therefore from slideN.xml naming assumptions).
        slide_id_list = prs.slides._sldIdLst
        elements = list(slide_id_list)
        for element in elements:
            slide_id_list.remove(element)
        for element in (elements[2], elements[1], elements[0]):
            slide_id_list.append(element)

        prs.save(path)
        return {
            "ids": ids,
            "layouts": layouts,
            "canvas": {
                "width": int(Inches(12)),
                "height": int(Inches(7)),
                "unit": "EMU",
            },
        }

    def _extract(self, name="output", visible_only=False):
        output_dir = self.root / name
        return EXTRACT.extract_pptx(
            self.pptx_path,
            output_dir,
            visible_only=visible_only,
        )

    def assert_shape_metadata(self, item, name, z_order, z_order_path):
        self.assertEqual(item["shape_name"], name)
        self.assertIsInstance(item["shape_id"], int)
        self.assertIsInstance(item["shape_type"], str)
        self.assertEqual(item["z_order"], z_order)
        self.assertEqual(item["z_order_path"], z_order_path)
        self.assertEqual(
            set(item["bbox"]), {"left", "top", "width", "height"}
        )
        self.assertTrue(all(isinstance(v, int) for v in item["bbox"].values()))

    def test_hidden_slide_id_layout_canvas_and_true_order(self):
        slides = self._extract()

        self.assertEqual([slide["number"] for slide in slides], [1, 2, 3])
        self.assertEqual(
            [slide["slide_id"] for slide in slides],
            [
                self.expected["ids"]["geometry"],
                self.expected["ids"]["hidden"],
                self.expected["ids"]["first"],
            ],
        )
        self.assertEqual(
            [slide["hidden"] for slide in slides], [False, True, False]
        )
        self.assertEqual(
            [slide["title"] for slide in slides],
            ["", "Hidden title", "First title"],
        )
        for slide in slides:
            self.assertTrue(
                {"number", "title", "content", "images", "notes"}
                <= set(slide)
            )
            self.assertEqual(slide["canvas"], self.expected["canvas"])
            self.assertEqual(
                slide["layout"], self.expected["layouts"][slide["slide_id"]]
            )

    def test_title_and_body_remain_separate_and_have_shape_metadata(self):
        slides = self._extract()
        first = next(
            slide
            for slide in slides
            if slide["slide_id"] == self.expected["ids"]["first"]
        )

        self.assertEqual(first["title"], "First title")
        self.assert_shape_metadata(
            first["title_shape"], "First Title Shape", 0, [0]
        )
        body_items = [
            item
            for item in first["content"]
            if item["type"] == "text" and item["content"] == "First body"
        ]
        self.assertEqual(len(body_items), 1)
        self.assert_shape_metadata(
            body_items[0], "First Body Shape", 1, [1]
        )

    def test_text_table_image_and_group_positions(self):
        geometry = self._extract()[0]
        content_by_name = {
            item["shape_name"]: item for item in geometry["content"]
        }
        images_by_name = {
            item["shape_name"]: item for item in geometry["images"]
        }

        text_item = content_by_name["Positioned Text"]
        self.assert_shape_metadata(text_item, "Positioned Text", 0, [0])
        self.assertEqual(text_item["shape_type"], "TEXT_BOX")
        self.assertEqual(
            text_item["bbox"],
            {
                "left": int(Inches(1)),
                "top": int(Inches(1)),
                "width": int(Inches(2)),
                "height": int(Inches(1)),
            },
        )

        table_item = content_by_name["Positioned Table"]
        self.assertEqual(table_item["rows"], [["A", "B"]])
        self.assert_shape_metadata(table_item, "Positioned Table", 1, [1])
        self.assertEqual(table_item["shape_type"], "TABLE")
        self.assertEqual(
            table_item["bbox"],
            {
                "left": int(Inches(3)),
                "top": int(Inches(1)),
                "width": int(Inches(2)),
                "height": int(Inches(1)),
            },
        )

        picture = images_by_name["Positioned Picture"]
        self.assert_shape_metadata(picture, "Positioned Picture", 2, [2])
        self.assertEqual(picture["shape_type"], "PICTURE")
        self.assertEqual(
            {key: picture[key] for key in ("left", "top", "width", "height")},
            picture["bbox"],
        )
        self.assertEqual(
            picture["bbox"],
            {
                "left": int(Inches(6)),
                "top": int(Inches(1)),
                "width": int(Inches(2)),
                "height": int(Inches(1)),
            },
        )

        grouped_text = content_by_name["Grouped Text"]
        self.assert_shape_metadata(grouped_text, "Grouped Text", 3, [3, 0])
        self.assertEqual(
            grouped_text["bbox"],
            {
                "left": int(Inches(7)),
                "top": int(Inches(4)),
                "width": int(Inches(2)),
                "height": int(Inches(2)),
            },
        )

        grouped_picture = images_by_name["Grouped Picture"]
        self.assert_shape_metadata(
            grouped_picture, "Grouped Picture", 3, [3, 1]
        )
        self.assertEqual(
            {key: grouped_picture[key] for key in ("left", "top", "width", "height")},
            grouped_picture["bbox"],
        )
        self.assertEqual(
            grouped_picture["bbox"],
            {
                "left": int(Inches(10)),
                "top": int(Inches(4)),
                "width": int(Inches(1)),
                "height": int(Inches(2)),
            },
        )

    def test_visible_only_api_preserves_source_numbers(self):
        slides = self._extract(visible_only=True)

        self.assertEqual([slide["number"] for slide in slides], [1, 3])
        self.assertTrue(all(not slide["hidden"] for slide in slides))
        self.assertEqual(
            [slide["slide_id"] for slide in slides],
            [
                self.expected["ids"]["geometry"],
                self.expected["ids"]["first"],
            ],
        )

    def test_cli_reports_hidden_slides_and_visible_only_filters_them(self):
        default_output = self.root / "cli-default"
        default_result = subprocess.run(
            [
                sys.executable,
                "-X",
                "utf8",
                str(SCRIPT),
                str(self.pptx_path),
                str(default_output),
            ],
            capture_output=True,
            text=True,
            encoding="utf-8",
        )
        self.assertEqual(default_result.returncode, 0, default_result.stderr)
        self.assertIn("Source-material extraction only", default_result.stdout)
        self.assertIn(
            "Hidden slides retained: 1 (source slide 2)",
            default_result.stdout,
        )
        default_data = json.loads(
            (default_output / "extracted-slides.json").read_text(encoding="utf-8")
        )
        self.assertEqual(len(default_data), 3)
        self.assertTrue(default_data[1]["hidden"])

        visible_output = self.root / "cli-visible"
        visible_result = subprocess.run(
            [
                sys.executable,
                "-X",
                "utf8",
                str(SCRIPT),
                str(self.pptx_path),
                str(visible_output),
                "--visible-only",
            ],
            capture_output=True,
            text=True,
            encoding="utf-8",
        )
        self.assertEqual(visible_result.returncode, 0, visible_result.stderr)
        self.assertIn(
            "Hidden slides omitted by --visible-only: 1 (source slide 2)",
            visible_result.stdout,
        )
        visible_data = json.loads(
            (visible_output / "extracted-slides.json").read_text(encoding="utf-8")
        )
        self.assertEqual([slide["number"] for slide in visible_data], [1, 3])
        self.assertTrue(all(not slide["hidden"] for slide in visible_data))


if __name__ == "__main__":
    unittest.main()
